#!/usr/bin/env python3
"""Сборка PowerPoint из структуры ПРЕЗЕНТАЦИЯ.md."""

from __future__ import annotations

import base64
import sys
from pathlib import Path

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
ASSETS = ROOT / "assets" / "presentation"
OUTPUT = ROOT / "assets" / "presentation" / "Mektep_Analyzer.pptx"

# Цвета бренда
C_PRIMARY = RGBColor(0x19, 0x76, 0xD2)  # синий
C_ACCENT = RGBColor(0x38, 0x8E, 0x3C)   # зелёный
C_DARK = RGBColor(0x26, 0x32, 0x38)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED = RGBColor(0x5F, 0x6B, 0x7A)
C_BG_DARK = RGBColor(0x1A, 0x23, 0x32)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def mermaid_png(mermaid_code: str, name: str) -> Path | None:
    """Скачать PNG диаграммы (Kroki → mermaid.ink fallback)."""
    cache = ASSETS / "mermaid"
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / f"{name}.png"
    if out.exists() and out.stat().st_size > 1000:
        return out

    sources = [
        ("kroki", lambda: requests.post(
            "https://kroki.io/mermaid/png",
            data=mermaid_code.strip().encode("utf-8"),
            headers={"Content-Type": "text/plain"},
            timeout=90,
        )),
        ("mermaid.ink-b64", lambda: requests.get(
            "https://mermaid.ink/img/"
            + base64.urlsafe_b64encode(mermaid_code.strip().encode("utf-8")).decode("ascii"),
            timeout=90,
        )),
    ]
    for label, fetch in sources:
        try:
            resp = fetch()
            resp.raise_for_status()
            if resp.content[:4] == b"\x89PNG":
                out.write_bytes(resp.content)
                print(f"  [ok] {name} via {label}")
                return out
        except Exception as exc:  # noqa: BLE001
            print(f"  [!] {name} ({label}): {exc}", file=sys.stderr)
    return None


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_overlay(slide, opacity: float = 0.45) -> None:
    """Полупрозрачная плашка под текст (имитация затемнения видео-фона)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        SLIDE_W,
        SLIDE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x32)
    shape.fill.transparency = 1.0 - opacity
    shape.line.fill.background()
    # отправить на задний план
    sp_tree = slide.shapes._spTree  # noqa: SLF001
    el = shape._element  # noqa: SLF001
    sp_tree.remove(el)
    sp_tree.insert(2, el)


def add_video_bg(slide, video_path: Path) -> bool:
    if not video_path.exists():
        return False
    slide.shapes.add_movie(
        str(video_path),
        Inches(0),
        Inches(0),
        SLIDE_W,
        SLIDE_H,
        poster_frame_image=None,
        mime_type="video/mp4",
    )
    return True


def textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = C_DARK,
    align=PP_ALIGN.LEFT,
    font_name: str = "Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font_name
    p.font.color.rgb = color
    p.alignment = align
    return box


def bullets(slide, left, top, width, height, items: list[str], *, size=14, color=C_DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # убрать **markdown**
        clean = item.replace("**", "")
        p.text = clean
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def table_slide(slide, left, top, width, headers: list[str], rows: list[list[str]], *, font_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.35 * n_rows))
    tbl = shape.table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.color.rgb = C_WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val.replace("**", "")
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size - 1)
                p.font.color.rgb = C_DARK
    return shape


def add_image_if_exists(slide, path: Path, left, top, width) -> bool:
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width)
        return True
    return False


def speaker_notes(slide, text: str) -> None:
    if not text.strip():
        return
    notes = slide.notes_slide.notes_text_frame
    notes.text = text.strip()


MERMAID_WAS = """flowchart TB
    MK[mektep.edu.kz] --> A1[выписывают оценки]
    A1 --> T[Учителя]
    T --> A2[заполняют Excel вручную]
    A2 --> R[5 видов отчётов]
    R --> A3[отправляют завучу]
    A3 --> Z[Завуч]
    Z --> A4[перенос в 1С]
    A4 --> F[Конечный отчёт]"""

MERMAID_NOW = """flowchart TB
    MK["mektep.edu.kz"] --> A1["извлечение оценок"]
    A1 --> D["Десктоп-приложение"]
    D --> R["5 видов отчётов"]
    D --> M["Мониторинг / обмен / своды"]
    GOALS["Цели учителя"] --> AI["ИИ-анализ СОР/СОЧ"] --> R
    R --> A2["загрузка на сервер"]
    A2 --> SRV["Аналитика + ИИ"]
    SRV --> WEB["Веб-кабинет завуча"]"""

MERMAID_EVOL = """flowchart LR
    subgraph B["БЫЛО"]
        BT["Учитель"] --> BZ["Завуч"] --> BO["1С"]
    end
    subgraph N["СТАЛО"]
        NK["mektep.edu.kz"] --> NP["Платформа"] --> NZ["Завуч"]
    end
    subgraph F["БУДЕТ"]
        FS["Школы"] --> FG["ГорОНО"] --> FO["ОблОНО"]
    end
    B --> N --> F"""

MERMAID_ROADMAP = """flowchart TB
    T["Учителя"] --> A1["отчёты"]
    A1 --> S["Школы"]
    S --> A2["ИИ-обобщение"]
    A2 --> GOR["ГорОНО"]
    GOR --> OBL["ОблОНО"]"""


def build() -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    print("Рендер Mermaid-диаграмм...")
    img_was = mermaid_png(MERMAID_WAS, "slide1-was")
    img_now = mermaid_png(MERMAID_NOW, "slide2-now")
    img_evol = mermaid_png(MERMAID_EVOL, "slide3-evol")
    img_road = mermaid_png(MERMAID_ROADMAP, "slide5-roadmap")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # blank

    demo_video = ASSETS / "demo-loop.mp4"

    # --- Слайд 0: Титул ---
    slide = prs.slides.add_slide(blank)
    if add_video_bg(slide, demo_video):
        add_overlay(slide, 0.5)
    else:
        set_slide_bg(slide, C_BG_DARK)
    textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
            "Mektep Analyzer", size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.8),
            "Платформа для автоматизации школьной отчётности",
            size=24, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6),
            "От ручного переписывания оценок — к готовому отчёту за минуты",
            size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
    if not demo_video.exists():
        textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                "💡 Добавьте demo-loop.mp4 в assets/presentation/ для видео-фона",
                size=12, color=C_MUTED, align=PP_ALIGN.CENTER)
    speaker_notes(slide, "Добрый день. Mektep Analyzer — платформа для автоматизации школьной отчётности. "
                  "На фоне — как это работает: сбор с mektep.edu.kz, отчёты учителю, аналитика завучу.")

    # --- Слайд 1: Как было ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, C_WHITE)
    textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
            "Как было (до платформы)", size=28, bold=True, color=C_PRIMARY)
    textbox(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
            "Вся отчётность — вручную, на каждом уровне", size=16, color=C_MUTED)
    if img_was:
        slide.shapes.add_picture(str(img_was), Inches(0.4), Inches(1.3), width=Inches(8.5))
    textbox(slide, Inches(9.1), Inches(1.5), Inches(3.8), Inches(2.5),
            "Итог: одни и те же данные переписываются руками дважды — "
            "учителем и завучом. Долго и с ошибками.",
            size=14, bold=True, color=RGBColor(0xC6, 0x28, 0x28))
    speaker_notes(slide, "Все данные на mektep.edu.kz. Учитель вручную выписывает оценки, "
                  "заполняет 5 отчётов, отправляет завучу. Завуч переносит в 1С. Двойной ввод.")

    # --- Слайд 2: Как стало ---
    slide = prs.slides.add_slide(blank)
    if add_video_bg(slide, demo_video):
        add_overlay(slide, 0.55)
        title_color, sub_color, bullet_color = C_WHITE, RGBColor(0xBB, 0xDE, 0xFB), C_WHITE
    else:
        set_slide_bg(slide, C_WHITE)
        title_color, sub_color, bullet_color = C_PRIMARY, C_MUTED, C_DARK
    textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
            "Как стало (сейчас)", size=28, bold=True, color=title_color)
    textbox(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
            "mektep.edu.kz → десктоп → сервер → завуч скачивает готовое",
            size=15, color=sub_color)
    if img_now and not demo_video.exists():
        slide.shapes.add_picture(str(img_now), Inches(0.3), Inches(1.2), width=Inches(6.5))
    bullets(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5), [
        "Ручной ввод учителя → 0: данные с mektep.edu.kz автоматически",
        "Десктоп: мониторинг, обмен данными, сводные таблицы",
        "Учитель задаёт только цели — ИИ пишет анализ СОР/СОЧ",
        "Ручное сведение завуча → 0: отчёты в веб-кабинете",
        "Сбор на компьютере учителя — пароли не уходят на сервер",
    ], size=15, color=bullet_color)
    speaker_notes(slide, "Данные извлекаются на компьютер учителя. Десктоп формирует отчёты, "
                  "ИИ пишет анализ. Завуч скачивает готовое из веб-кабинета.")

    # --- Слайд 2.5: Демо ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, C_WHITE)
    textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
            "Платформа в работе", size=28, bold=True, color=C_PRIMARY)
    textbox(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
            "От портала — к готовому отчёту за минуты", size=16, color=C_MUTED)
    table_slide(
        slide, Inches(0.5), Inches(1.4), Inches(12),
        ["Шаг", "Что видит зал"],
        [
            ["1", "Десктоп: автосбор, прогресс в реальном времени"],
            ["2", "Формирование отчёта, список готовых файлов"],
            ["3", "Веб-кабинет завуча: аналитика, графики, скачивание"],
        ],
    )
    y_img = Inches(3.2)
    w_img = Inches(3.8)
    imgs = [
        (ASSETS / "01-desktop-scrape.png", Inches(0.5)),
        (ASSETS / "02-desktop-reports.png", Inches(4.6)),
        (ASSETS / "03-admin-dashboard.png", Inches(8.7)),
    ]
    any_img = False
    for path, left in imgs:
        if add_image_if_exists(slide, path, left, y_img, w_img):
            any_img = True
    if not any_img:
        # плейсхолдеры
        for idx, (_, left) in enumerate(imgs):
            shape = slide.shapes.add_shape(1, left, y_img, w_img, Inches(2.8))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xEC, 0xEF, 0xF1)
            shape.line.color.rgb = C_MUTED
            textbox(slide, left, y_img + Inches(1.0), w_img, Inches(0.8),
                    f"Скрин {idx + 1}\n(добавьте PNG в assets/presentation/)",
                    size=11, color=C_MUTED, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
            "45 сек — показать продукт «вживую» или demo-loop.mp4",
            size=12, color=C_MUTED)
    speaker_notes(slide, "Учитель: Собрать данные → отчёты в один клик → завуч: веб-кабинет.")

    # --- Слайд 3: Эволюция ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, C_WHITE)
    textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
            "Эволюция: кто перестаёт работать руками", size=26, bold=True, color=C_PRIMARY)
    table_slide(
        slide, Inches(0.5), Inches(0.95), Inches(12),
        ["Этап", "Кто работал вручную", "Что убрала платформа"],
        [
            ["Было", "Учитель + Завуч", "—"],
            ["Стало", "Никто (учитель задаёт только цели)",
             "Ввод учителя, сведение завуча, текст анализа СОР/СОЧ"],
            ["Будет", "Никто", "Итоговое обобщение + сводки ГорОНО / ОблОНО"],
        ],
        font_size=10,
    )
    if img_evol:
        slide.shapes.add_picture(str(img_evol), Inches(0.5), Inches(3.0), width=Inches(12))
    speaker_notes(slide, "Было: учитель и завуч. Стало: платформа. Будет: ГорОНО и ОблОНО.")

    # --- Слайд 4: Возможности ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, C_WHITE)
    textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
            "Что умеет платформа сегодня", size=28, bold=True, color=C_PRIMARY)
    table_slide(
        slide, Inches(0.4), Inches(0.95), Inches(7.5),
        ["Возможность", "Описание"],
        [
            ["Автосбор данных", "СОР/СОЧ, прогресс в реальном времени"],
            ["Десктоп", "Мониторинг, обмен, сводные таблицы"],
            ["Отчёты в один клик", "Excel + Word, рус./каз."],
            ["ИИ-анализ СОР/СОЧ", "Затруднения, причины, план работы"],
            ["Аналитика для завуча", "Качество, успеваемость, графики"],
        ],
        font_size=10,
    )
    collage = ASSETS / "demo-collage.png"
    if not add_image_if_exists(slide, collage, Inches(8.1), Inches(0.95), Inches(4.8)):
        shape = slide.shapes.add_shape(1, Inches(8.1), Inches(0.95), Inches(4.8), Inches(5.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xEC, 0xEF, 0xF1)
        textbox(slide, Inches(8.1), Inches(3.0), Inches(4.8), Inches(1),
                "demo-collage.png\n(скрины интерфейса)", size=11, color=C_MUTED, align=PP_ALIGN.CENTER)
    speaker_notes(slide, "Автосбор, десктоп, отчёты, ИИ, аналитика завуча.")

    # --- Слайд 5: Дорожная карта ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, C_WHITE)
    textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
            "Дальнейшее развитие", size=28, bold=True, color=C_PRIMARY)
    textbox(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
            "Дорожная карта: от школы к уровню области", size=16, color=C_MUTED)
    bullets(slide, Inches(0.5), Inches(1.35), Inches(5.5), Inches(3.5), [
        "1. ИИ для обобщения конечных отчётов — выводы, тенденции, проблемные зоны",
        "2. Админка ГорОНО / ОблОНО — только чтение и скачивание",
        "3. Конструктор таблиц для ГорОНО — свод по предметам и параллелям",
    ], size=13)
    if img_road:
        slide.shapes.add_picture(str(img_road), Inches(6.2), Inches(1.2), width=Inches(6.5))
    textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
            "Программа выходит за пределы школы и работает на уровне области.",
            size=14, bold=True, color=C_ACCENT)
    speaker_notes(slide, "ИИ-обобщение, многоуровневая админка, конструктор таблиц.")

    # --- Слайд 6: Эффект ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, C_WHITE)
    textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
            "Что получают школа и система образования", size=26, bold=True, color=C_PRIMARY)
    bullets(slide, Inches(0.7), Inches(1.1), Inches(11), Inches(4.5), [
        "Экономия времени — часы → минуты, без двойного ввода",
        "Меньше ошибок — расчёты делает система",
        "Единый стандарт отчётов на всех уровнях",
        "Наглядная аналитика успеваемости в реальном времени",
        "Двуязычность — русский и казахский",
        "Масштаб — от одной школы до района и области",
    ], size=18)
    textbox(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
            "Главный эффект: педагоги занимаются обучением, "
            "а управленцы — решениями, а не таблицами.",
            size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    speaker_notes(slide, "Экономия времени, меньше ошибок, единый стандарт, аналитика, масштаб.")

    # --- Слайд 7: Финал ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, C_BG_DARK)
    textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1),
            "Mektep Analyzer", size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.7),
            "Меньше рутины — больше времени на учеников",
            size=24, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
    textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
            "Контакты · ссылка на скачивание · QR-код демо",
            size=16, color=C_MUTED, align=PP_ALIGN.CENTER)
    speaker_notes(slide, "Спасибо за внимание. Готов ответить на вопросы и показать систему.")

    prs.save(str(OUTPUT))
    print(f"Готово: {OUTPUT}")
    return OUTPUT


if __name__ == "__main__":
    build()
